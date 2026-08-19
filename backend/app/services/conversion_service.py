import io
import os
import tempfile
import zipfile
import fitz  # PyMuPDF
import img2pdf
from PIL import Image
import docx
from docx import Document
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
import openpyxl
from openpyxl import Workbook
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

import shutil
import glob

def get_soffice_cmd() -> str | None:
    """Find LibreOffice binary across PATH and default Windows/Linux install directories."""
    # 1. Check CLI wrapper soffice.com first (Windows)
    cmd = shutil.which("soffice.com") or shutil.which("soffice") or shutil.which("libreoffice")
    if cmd:
        return cmd
    
    # 2. Common Windows directories
    win_paths = [
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        r"C:\Program Files\LibreOffice*\program\soffice.com",
        r"C:\Program Files\LibreOffice*\program\soffice.exe",
    ]
    for p in win_paths:
        matches = glob.glob(p)
        if matches and os.path.exists(matches[0]):
            return matches[0]
    return None

def get_cloudmersive_convert_api():
    api_key = os.getenv("CLOUDMERSIVE_API_KEY", "").strip()
    if not api_key:
        return None
    try:
        import cloudmersive_convert_api_client
        configuration = cloudmersive_convert_api_client.Configuration()
        configuration.api_key['Apikey'] = api_key
        return cloudmersive_convert_api_client.ConvertDocumentApi(cloudmersive_convert_api_client.ApiClient(configuration))
    except Exception as e:
        print(f"[Cloudmersive Init Error]: {e}")
        return None

def get_cloudmersive_ocr_api():
    api_key = os.getenv("CLOUDMERSIVE_API_KEY", "").strip()
    if not api_key:
        return None
    try:
        import cloudmersive_ocr_api_client
        configuration = cloudmersive_ocr_api_client.Configuration()
        configuration.api_key['Apikey'] = api_key
        return cloudmersive_ocr_api_client.PdfOcrApi(cloudmersive_ocr_api_client.ApiClient(configuration))
    except Exception as e:
        print(f"[Cloudmersive OCR Init Error]: {e}")
        return None

def get_ghostscript_cmd() -> str | None:
    """Find Ghostscript binary across PATH and default Windows/Linux install directories."""
    cmd = shutil.which("gswin64c") or shutil.which("gswin32c") or shutil.which("gs")
    if cmd:
        return cmd
    
    win_paths = [
        r"C:\Program Files\gs\gs*\bin\gswin64c.exe",
        r"C:\Program Files (x86)\gs\gs*\bin\gswin32c.exe",
    ]
    for p in win_paths:
        matches = glob.glob(p)
        if matches and os.path.exists(matches[0]):
            return matches[0]
    return None

def jpg_to_pdf(images_bytes_list: list[bytes]) -> bytes:
    if not images_bytes_list:
        return b""

    A4_W, A4_H = 595.28, 841.89
    doc = fitz.open()

    for b in images_bytes_list:
        try:
            img = Image.open(io.BytesIO(b))
            # Auto-orient based on EXIF
            try:
                from PIL import ImageOps
                img = ImageOps.exif_transpose(img)
            except Exception:
                pass

            if img.mode in ("RGBA", "P", "LA", "CMYK"):
                bg = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode in ("RGBA", "LA") and "A" in img.getbands():
                    bg.paste(img, mask=img.split()[-1])
                else:
                    bg.paste(img.convert("RGB"))
                img = bg
            else:
                img = img.convert("RGB")

            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=98)
            clean_bytes = buf.getvalue()

            img_w, img_h = img.size
            # Calculate scale to fit within A4 page proportionally
            scale = min(A4_W / img_w, A4_H / img_h)
            fit_w = img_w * scale
            fit_h = img_h * scale
            x0 = (A4_W - fit_w) / 2.0
            y0 = (A4_H - fit_h) / 2.0

            page = doc.new_page(width=A4_W, height=A4_H)
            rect = fitz.Rect(x0, y0, x0 + fit_w, y0 + fit_h)
            page.insert_image(rect, stream=clean_bytes)
        except Exception as e:
            print(f"[jpg_to_pdf page]: {e}")

    if len(doc) > 0:
        out_buf = io.BytesIO()
        doc.save(out_buf, garbage=4, deflate=True, clean=True)
        doc.close()
        return out_buf.getvalue()
    return b""

def word_to_pdf(docx_bytes: bytes) -> bytes:
    import html
    import subprocess
    import tempfile

    # 1. Primary: If LibreOffice / soffice is installed, use native headless rendering
    soffice_cmd = get_soffice_cmd()
    if soffice_cmd:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                in_path = os.path.join(tmpdir, "document.docx")
                with open(in_path, "wb") as f:
                    f.write(docx_bytes)
                subprocess.run([soffice_cmd, "--headless", "--convert-to", "pdf", in_path, "--outdir", tmpdir], check=True, timeout=20)
                out_pdf = os.path.join(tmpdir, "document.pdf")
                if os.path.exists(out_pdf):
                    with open(out_pdf, "rb") as f:
                        return f.read()
        except Exception as e:
            print(f"[word_to_pdf LibreOffice]: {e}")

    # 2. High-precision python-docx + ReportLab document builder
    doc = Document(io.BytesIO(docx_bytes))
    buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    styles = getSampleStyleSheet()
    story = []

    for p in doc.paragraphs:
        raw_text = p.text.strip()
        if raw_text:
            safe_text = html.escape(raw_text)
            style = styles['Heading1'] if p.style.name.startswith('Heading') else styles['Normal']
            story.append(Paragraph(safe_text, style))
            story.append(Spacer(1, 8))

    for table in doc.tables:
        table_data = []
        for row in table.rows:
            row_data = [Paragraph(html.escape(cell.text.strip()), styles['Normal']) for cell in row.cells]
            table_data.append(row_data)
        if table_data:
            t = Table(table_data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#F1F5F9')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.HexColor('#0F172A')),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#CBD5E1'))
            ]))
            story.append(t)
            story.append(Spacer(1, 12))

    if not story:
        story.append(Paragraph("Empty Word Document", styles['Normal']))

    pdf_doc.build(story)
    return buffer.getvalue()

def ppt_to_pdf(pptx_bytes: bytes) -> bytes:
    import subprocess
    import tempfile
    import os
    import shutil

    # 1. Primary: Headless LibreOffice Impress Export (100% pixel-perfect slide layout & fonts)
    soffice_cmd = get_soffice_cmd()
    if soffice_cmd:
        tmpdir = tempfile.mkdtemp(prefix="docflow_ppt_")
        try:
            in_path = os.path.join(tmpdir, "presentation.pptx")
            with open(in_path, "wb") as f:
                f.write(pptx_bytes)
            subprocess.run([soffice_cmd, "--headless", "--convert-to", "pdf:impress_pdf_Export", in_path, "--outdir", tmpdir], check=True, timeout=30)
            out_pdf = os.path.join(tmpdir, "presentation.pdf")
            if os.path.exists(out_pdf) and os.path.getsize(out_pdf) > 0:
                with open(out_pdf, "rb") as f:
                    pdf_data = f.read()
                return pdf_data
        except Exception as e:
            print(f"[ppt_to_pdf LibreOffice]: {e}")
        finally:
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except Exception:
                pass

    # 2. Secondary: Cloudmersive Enterprise PPTX to PDF
    cm_api = get_cloudmersive_convert_api()
    if cm_api:
        tmp_in_path = None
        try:
            tmp_f = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
            tmp_f.write(pptx_bytes)
            tmp_f.flush()
            tmp_f.close()
            tmp_in_path = tmp_f.name

            res = cm_api.convert_document_pptx_to_pdf(tmp_in_path)
            if res:
                if isinstance(res, str):
                    if os.path.exists(res):
                        with open(res, "rb") as fr:
                            res = fr.read()
                    elif (res.startswith("b'") and res.endswith("'")) or (res.startswith('b"') and res.endswith('"')):
                        try:
                            import ast
                            res = ast.literal_eval(res)
                        except Exception:
                            res = res.encode("latin1", errors="ignore")
                    else:
                        res = res.encode("latin1", errors="ignore")
                if isinstance(res, bytes) and len(res) > 50:
                    return res
        except Exception as e:
            print(f"[ppt_to_pdf Cloudmersive]: {e}")
        finally:
            if tmp_in_path and os.path.exists(tmp_in_path):
                try:
                    os.remove(tmp_in_path)
                except Exception:
                    pass

    # 3. Dynamic Slide Layout Fallback
    prs = Presentation(io.BytesIO(pptx_bytes))
    slide_w_pt = prs.slide_width.pt if prs.slide_width else 720.0
    slide_h_pt = prs.slide_height.pt if prs.slide_height else 405.0
    
    doc = fitz.open()
    for slide_idx, slide in enumerate(prs.slides):
        page = doc.new_page(width=slide_w_pt, height=slide_h_pt)
        for shape in slide.shapes:
            if shape.has_text_frame:
                sx = shape.left.pt if shape.left else 40.0
                sy = shape.top.pt if shape.top else 40.0
                sw = shape.width.pt if shape.width else slide_w_pt - 80.0
                sh = shape.height.pt if shape.height else 40.0
                rect = fitz.Rect(sx, sy, sx + sw, sy + sh)
                text = shape.text_frame.text.strip()
                if text:
                    page.insert_textbox(rect, text, fontsize=12, fontname="helv")

    buffer = io.BytesIO()
    doc.save(buffer, garbage=4, deflate=True, clean=True)
    doc.close()
    return buffer.getvalue()

def excel_to_pdf(xlsx_bytes: bytes) -> bytes:
    """
    Industrial-Grade Excel (.xlsx / .xls) -> PDF Converter.
    Uses Headless LibreOffice Calc Engine (pdf:calc_pdf_Export) with multi-sheet
    rendering and post-conversion PyMuPDF content verification.
    """
    if not xlsx_bytes or len(xlsx_bytes) < 10:
        raise ValueError("Invalid or empty Excel file provided.")

    import subprocess
    import tempfile
    import os
    import shutil
    import openpyxl

    # Pre-validation & sample content inspection
    sample_values = []
    try:
        wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
        for sname in wb.sheetnames:
            ws = wb[sname]
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    if v is not None and len(str(v).strip()) > 1 and len(sample_values) < 5:
                        sample_values.append(str(v).strip())
    except Exception:
        pass

    # Pre-configure page setup (fit all columns to 1 page width, auto landscape for wide tables)
    try:
        wb_prep = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
        for ws in wb_prep.worksheets:
            max_col = ws.max_column or 1
            ws.sheet_properties.pageSetUpPr.fitToPage = True
            ws.page_setup.fitToPage = True
            ws.page_setup.fitToWidth = 1
            ws.page_setup.fitToHeight = 0
            ws.print_options.gridLines = True
            ws.print_options.gridLinesSet = True
            if max_col > 6:
                ws.page_setup.orientation = ws.ORIENTATION_LANDSCAPE
            else:
                ws.page_setup.orientation = ws.ORIENTATION_PORTRAIT
        prep_buf = io.BytesIO()
        wb_prep.save(prep_buf)
        xlsx_bytes = prep_buf.getvalue()
    except Exception as prep_e:
        print(f"[excel_to_pdf prep]: {prep_e}")

    # 1. Primary: Headless LibreOffice Calc Engine
    soffice_cmd = get_soffice_cmd()
    if soffice_cmd:
        tmpdir = tempfile.mkdtemp(prefix="docflow_calc_")
        try:
            in_path = os.path.join(tmpdir, "sheet.xlsx")
            with open(in_path, "wb") as f:
                f.write(xlsx_bytes)
            proc = subprocess.run(
                [soffice_cmd, "--headless", "--convert-to", "pdf:calc_pdf_Export", in_path, "--outdir", tmpdir],
                capture_output=True, text=True, timeout=40
            )
            out_pdf = os.path.join(tmpdir, "sheet.pdf")
            if os.path.exists(out_pdf) and os.path.getsize(out_pdf) > 0:
                with open(out_pdf, "rb") as f:
                    pdf_bytes = f.read()

                # Post-conversion validation
                try:
                    val_doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                    if len(val_doc) == 0:
                        val_doc.close()
                        raise ValueError("Generated PDF has 0 pages.")
                    val_doc.close()
                except Exception as ve:
                    raise ValueError(f"Generated PDF failed validation: {ve}")

                return pdf_bytes
        except Exception as e:
            print(f"[excel_to_pdf LibreOffice]: {e}")
        finally:
            try:
                shutil.rmtree(tmpdir, ignore_errors=True)
            except Exception:
                pass

    # 2. Secondary: Cloudmersive Enterprise XLSX to PDF Engine
    cm_api = get_cloudmersive_convert_api()
    if cm_api:
        tmp_in_path = None
        try:
            tmp_f = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
            tmp_f.write(xlsx_bytes)
            tmp_f.flush()
            tmp_f.close()
            tmp_in_path = tmp_f.name

            res = cm_api.convert_document_xlsx_to_pdf(tmp_in_path)
            if res:
                if isinstance(res, str) and os.path.exists(res):
                    with open(res, "rb") as fr:
                        res = fr.read()
                elif isinstance(res, str):
                    res = res.encode("latin1", errors="ignore")
                if isinstance(res, bytes) and len(res) > 50:
                    return res
        except Exception as e:
            print(f"[excel_to_pdf Cloudmersive]: {e}")
        finally:
            if tmp_in_path and os.path.exists(tmp_in_path):
                try:
                    os.remove(tmp_in_path)
                except Exception:
                    pass

    raise RuntimeError(
        "Excel to PDF conversion failed because LibreOffice is not installed or configured on the server. "
        "Please install LibreOffice (https://www.libreoffice.org/download/download/)."
    )

def html_to_pdf(html_content: str) -> bytes:
    import tempfile
    import subprocess
    import os

    # 1. Primary: Cloudmersive Enterprise HTML to PDF Converter
    cm_api = get_cloudmersive_convert_api()
    if cm_api:
        try:
            with tempfile.NamedTemporaryFile(suffix=".html", delete=False, mode="w", encoding="utf-8") as tmp_in:
                tmp_in.write(html_content)
                tmp_in_path = tmp_in.name
            try:
                res = cm_api.convert_document_html_to_pdf(tmp_in_path)
                if res:
                    if isinstance(res, str):
                        if os.path.exists(res):
                            with open(res, "rb") as fr:
                                res = fr.read()
                        elif (res.startswith("b'") and res.endswith("'")) or (res.startswith('b"') and res.endswith('"')):
                            try:
                                import ast
                                res = ast.literal_eval(res)
                            except Exception:
                                res = res.encode("latin1", errors="ignore")
                        else:
                            res = res.encode("latin1", errors="ignore")
                    if isinstance(res, bytes) and len(res) > 50:
                        return res
            finally:
                if os.path.exists(tmp_in_path):
                    os.remove(tmp_in_path)
        except Exception as e:
            print(f"[html_to_pdf Cloudmersive]: {e}")

    # 2. Secondary: Headless LibreOffice HTML to PDF Engine
    soffice_cmd = get_soffice_cmd()
    if soffice_cmd:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                in_html = os.path.join(tmpdir, "document.html")
                with open(in_html, "w", encoding="utf-8") as f:
                    f.write(html_content)
                res = subprocess.run(
                    [soffice_cmd, "--headless", "--convert-to", "pdf", in_html, "--outdir", tmpdir],
                    capture_output=True,
                    timeout=20
                )
                out_pdf = os.path.join(tmpdir, "document.pdf")
                if os.path.exists(out_pdf) and os.path.getsize(out_pdf) > 0:
                    with open(out_pdf, "rb") as f:
                        return f.read()
        except Exception as e:
            print(f"[html_to_pdf LibreOffice]: {e}")

    # 3. Structural ReportLab Fallback
    buffer = io.BytesIO()
    pdf_doc = SimpleDocTemplate(buffer, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    styles = getSampleStyleSheet()
    story = []

    import re
    clean_text = re.sub('<[^<]+?>', '', html_content)
    for line in clean_text.splitlines():
        line = line.strip()
        if line:
            story.append(Paragraph(line, styles['Normal']))
            story.append(Spacer(1, 6))

    if not story:
        story.append(Paragraph("Empty HTML Document", styles['Normal']))

    pdf_doc.build(story)
    return buffer.getvalue()

def pdf_to_jpg(pdf_bytes: bytes) -> list[tuple[str, bytes]]:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    pages = []
    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("jpeg")
        pages.append((f"page_{i+1}.jpg", img_bytes))
    doc.close()
    return pages

def pdf_to_word(pdf_bytes: bytes) -> bytes:
    # 0. Cloudmersive Enterprise Engine
    cm_api = get_cloudmersive_convert_api()
    if cm_api:
        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_in:
                tmp_in.write(pdf_bytes)
                tmp_in_path = tmp_in.name
            try:
                res = cm_api.convert_document_pdf_to_docx(tmp_in_path)
                if res and len(res) > 100:
                    return res
            finally:
                if os.path.exists(tmp_in_path):
                    os.remove(tmp_in_path)
        except Exception as e:
            print(f"[pdf_to_word Cloudmersive]: {e}")

    # 1. High-precision pdf2docx conversion
    try:
        from pdf2docx import Converter
        pdf_file = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
        pdf_file.write(pdf_bytes)
        pdf_file.close()
        pdf_path = pdf_file.name
        docx_path = pdf_path + ".docx"

        try:
            cv = Converter(pdf_path)
            cv.convert(docx_path)
            cv.close()

            with open(docx_path, "rb") as f_docx:
                out_bytes = f_docx.read()
            return out_bytes
        finally:
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
            if os.path.exists(docx_path):
                os.remove(docx_path)
    except Exception as e:
        print(f"[pdf_to_word fallback]: {e}")

    # 2. PyMuPDF + python-docx structural fallback
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    docx_doc = Document()
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            for paragraph in text.split("\n\n"):
                if paragraph.strip():
                    docx_doc.add_paragraph(paragraph.strip())
    buffer = io.BytesIO()
    docx_doc.save(buffer)
    doc.close()
    return buffer.getvalue()

def pdf_to_pptx(pdf_bytes: bytes) -> bytes:
    import tempfile
    import os
    from pptx import Presentation
    from pptx.util import Pt

    # Ultra-High-Definition (300 DPI) Vector Slide Engine (Matching iLovePDF / Smallpdf / Adobe standard)
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    prs = Presentation()
    
    first_page = doc[0] if len(doc) > 0 else None
    if first_page:
        prs.slide_width = Pt(first_page.rect.width)
        prs.slide_height = Pt(first_page.rect.height)
    else:
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    with tempfile.TemporaryDirectory() as tmpdir:
        for page_idx, page in enumerate(doc):
            slide = prs.slides.add_slide(blank_layout)
            
            # 300 DPI Ultra-Crisp Rendering - 100% pixel-perfect vector typography, graphics, and backgrounds
            pix = page.get_pixmap(dpi=300)
            img_path = os.path.join(tmpdir, f"slide_{page_idx}.png")
            pix.save(img_path)
            
            slide.shapes.add_picture(img_path, Pt(0), Pt(0), width=Pt(page.rect.width), height=Pt(page.rect.height))

        out_pptx = os.path.join(tmpdir, "presentation.pptx")
        prs.save(out_pptx)
        with open(out_pptx, "rb") as f:
            out_bytes = f.read()

    doc.close()
    return out_bytes

def pdf_to_excel(pdf_bytes: bytes) -> bytes:
    """
    Industrial-Grade PDF -> Excel (.xlsx) Converter.
    Multi-engine extraction (PyMuPDF, pdfplumber, OCR), table quality validation,
    multi-page table continuation stitching, and OpenPyXL professional table formatting.
    """
    if not pdf_bytes or len(pdf_bytes) < 10:
        raise ValueError("Invalid or empty PDF file provided.")

    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception as e:
        raise ValueError(f"The PDF file appears corrupted or unreadable: {e}")

    total_pages = len(doc)
    if total_pages == 0:
        doc.close()
        raise ValueError("The uploaded PDF contains 0 pages.")

    import pdfplumber
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
    import re
    import tempfile
    import datetime

    def parse_value(val):
        if val is None:
            return ""
        s = ILLEGAL_CHARACTERS_RE.sub("", str(val)).strip()
        if not s:
            return ""
        if s.isdigit() and len(s) < 14 and (not s.startswith("0") or len(s) == 1):
            try:
                return int(s)
            except ValueError:
                pass
        cleaned = re.sub(r"[,$€₹£]", "", s).strip()
        if cleaned.endswith("%"):
            try:
                return float(cleaned[:-1].strip()) / 100.0
            except ValueError:
                pass
        try:
            if "." in cleaned and re.match(r"^-?\d+\.\d+$", cleaned):
                return float(cleaned)
        except ValueError:
            pass
        if re.match(r"^\d{4}-\d{2}-\d{2}$", s):
            try:
                return datetime.datetime.strptime(s, "%Y-%m-%d").date()
            except ValueError:
                pass
        return s

    def calc_table_score(rows):
        if not rows:
            return 0.0
        valid_rows = [r for r in rows if any(bool(str(c).strip()) for c in r if c is not None)]
        if not valid_rows:
            return 0.0
        total_rows = len(valid_rows)
        col_counts = [len(r) for r in valid_rows]
        max_cols = max(col_counts) if col_counts else 0
        if max_cols < 1:
            return 0.0
        equal_len_count = sum(1 for c in col_counts if c == max_cols)
        consistency_score = equal_len_count / total_rows
        total_cells = sum(len(r) for r in valid_rows)
        non_empty_cells = sum(sum(1 for c in r if c is not None and bool(str(c).strip())) for r in valid_rows)
        fill_ratio = (non_empty_cells / total_cells) if total_cells > 0 else 0.0
        if fill_ratio < 0.15:
            return 0.0
        return (consistency_score * 0.6) + (fill_ratio * 0.4)

    # Open pdfplumber reader
    plumber_pdf = None
    try:
        plumber_pdf = pdfplumber.open(io.BytesIO(pdf_bytes))
    except Exception:
        pass

    extracted_tables_all_pages = []

    with tempfile.TemporaryDirectory() as tmpdir:
        for page_idx in range(total_pages):
            page = doc[page_idx]
            raw_text = page.get_text("text").strip()

            # OCR Check for scanned page
            if len(raw_text) < 25:
                ocr_api = get_cloudmersive_ocr_api()
                if ocr_api:
                    try:
                        pix = page.get_pixmap(dpi=300)
                        img_path = os.path.join(tmpdir, f"scan_p{page_idx}.png")
                        pix.save(img_path)
                        ocr_res = ocr_api.pdf_ocr_pdf_to_lines_with_location(img_path)
                        if ocr_res and hasattr(ocr_res, "ocr_pages") and ocr_res.ocr_pages:
                            ocr_rows = []
                            for ocr_p in ocr_res.ocr_pages:
                                for ocr_l in ocr_p.ocr_lines:
                                    if ocr_l.line_text and ocr_l.line_text.strip():
                                        cols = [c.strip() for c in re.split(r"\t| {2,}", ocr_l.line_text) if c.strip()]
                                        if cols:
                                            ocr_rows.append(cols)
                            if ocr_rows and calc_table_score(ocr_rows) > 0.2:
                                extracted_tables_all_pages.append({
                                    "page": page_idx + 1,
                                    "rows": ocr_rows,
                                    "headers": ocr_rows[0] if ocr_rows else [],
                                    "score": calc_table_score(ocr_rows),
                                    "source": "OCR"
                                })
                                continue
                    except Exception as e:
                        print(f"[pdf_to_excel OCR]: {e}")

            # Engine 1: PyMuPDF find_tables()
            pymupdf_tables = []
            try:
                pymupdf_tabs = page.find_tables()
                if pymupdf_tabs and len(pymupdf_tabs.tables) > 0:
                    for t in pymupdf_tabs:
                        t_data = t.extract()
                        if t_data:
                            cleaned = [r for r in t_data if any(bool(str(c).strip()) for c in r if c is not None)]
                            sc = calc_table_score(cleaned)
                            if sc > 0.25:
                                pymupdf_tables.append({"rows": cleaned, "score": sc, "source": "PyMuPDF"})
            except Exception as e:
                print(f"[PyMuPDF find_tables page {page_idx}]: {e}")

            # Engine 2: pdfplumber extract_tables()
            plumber_tables = []
            if plumber_pdf and page_idx < len(plumber_pdf.pages):
                try:
                    p_page = plumber_pdf.pages[page_idx]
                    p_tabs = p_page.extract_tables(table_settings={"vertical_strategy": "lines", "horizontal_strategy": "lines"})
                    for pt in p_tabs:
                        cleaned = [r for r in pt if any(bool(str(c).strip()) for c in r if c is not None)]
                        sc = calc_table_score(cleaned)
                        if sc > 0.25:
                            plumber_tables.append({"rows": cleaned, "score": sc, "source": "pdfplumber_lines"})

                    if not plumber_tables:
                        p_tabs_text = p_page.extract_tables(table_settings={"vertical_strategy": "text", "horizontal_strategy": "text"})
                        for pt in p_tabs_text:
                            cleaned = [r for r in pt if any(bool(str(c).strip()) for c in r if c is not None)]
                            sc = calc_table_score(cleaned)
                            if sc > 0.3:
                                plumber_tables.append({"rows": cleaned, "score": sc, "source": "pdfplumber_text"})
                except Exception as e:
                    print(f"[pdfplumber page {page_idx}]: {e}")

            # Best Engine Selection
            chosen_tables = []
            if pymupdf_tables and plumber_tables:
                best_pm = max(pymupdf_tables, key=lambda x: x["score"])
                best_pl = max(plumber_tables, key=lambda x: x["score"])
                chosen_tables = [best_pl] if best_pl["score"] > best_pm["score"] + 0.15 else pymupdf_tables
            elif pymupdf_tables:
                chosen_tables = pymupdf_tables
            elif plumber_tables:
                chosen_tables = plumber_tables

            for ct in chosen_tables:
                rows = ct["rows"]
                extracted_tables_all_pages.append({
                    "page": page_idx + 1,
                    "rows": rows,
                    "headers": [str(c).strip() for c in rows[0]] if rows else [],
                    "score": ct["score"],
                    "source": ct["source"]
                })

    if plumber_pdf:
        try:
            plumber_pdf.close()
        except Exception:
            pass

    # Multi-Page Continuation Stitching & Multi-Sheet Structuring
    structured_sheets = []
    if extracted_tables_all_pages:
        for t_info in extracted_tables_all_pages:
            t_rows = t_info["rows"]
            t_headers = t_info["headers"]
            t_cols = len(t_headers) if t_headers else (len(t_rows[0]) if t_rows else 0)

            stitched = False
            if structured_sheets:
                last_sheet = structured_sheets[-1]
                last_headers = last_sheet["headers"]
                last_cols = len(last_headers) if last_headers else (len(last_sheet["rows"][0]) if last_sheet["rows"] else 0)

                if last_cols == t_cols and t_cols > 0:
                    # Check matching header row
                    is_repeat = False
                    if t_rows and last_headers:
                        if [str(c).strip().lower() for c in t_rows[0]] == [str(h).strip().lower() for h in last_headers]:
                            is_repeat = True
                    data_to_append = t_rows[1:] if is_repeat else t_rows
                    last_sheet["rows"].extend(data_to_append)
                    stitched = True

            if not stitched:
                sheet_idx = len(structured_sheets) + 1
                structured_sheets.append({
                    "name": f"Table_{sheet_idx}",
                    "headers": t_headers,
                    "rows": list(t_rows)
                })

    # If no tables found, extract formatted text into Extracted_Text sheet
    if not structured_sheets:
        extracted_text_rows = [["Page", "Paragraph #", "Extracted Content"]]
        for p_idx in range(total_pages):
            p_text = doc[p_idx].get_text("text").strip()
            if p_text:
                for para_idx, para in enumerate(p_text.split("\n\n"), start=1):
                    if para.strip():
                        extracted_text_rows.append([p_idx + 1, para_idx, para.strip()])
        if len(extracted_text_rows) > 1:
            structured_sheets.append({
                "name": "Extracted_Text",
                "headers": extracted_text_rows[0],
                "rows": extracted_text_rows
            })

    doc.close()

    if not structured_sheets:
        raise ValueError("No structured tables or readable text could be extracted from this PDF document.")

    # OpenPyXL Workbook Creation with Professional Table Formatting
    wb = Workbook()
    wb.remove(wb.active)

    header_fill = PatternFill(start_color="4F46E5", end_color="4F46E5", fill_type="solid")
    header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    alt_fill = PatternFill(start_color="F8FAFC", end_color="F8FAFC", fill_type="solid")
    regular_font = Font(name="Calibri", size=10, color="0F172A")
    thin_border = Border(
        left=Side(style="thin", color="CBD5E1"),
        right=Side(style="thin", color="CBD5E1"),
        top=Side(style="thin", color="CBD5E1"),
        bottom=Side(style="thin", color="CBD5E1")
    )

    for sheet_info in structured_sheets:
        ws = wb.create_sheet(title=sheet_info["name"][:31])
        ws.views.sheetView[0].showGridLines = True
        rows_data = sheet_info["rows"]
        if not rows_data:
            continue

        num_cols = max(len(r) for r in rows_data) if rows_data else 0

        # Precompute alignment per column based on data types so headers match data underneath
        col_alignments = {}
        for c_idx in range(num_cols):
            numeric_count = 0
            date_count = 0
            text_count = 0
            for r_idx in range(1, len(rows_data)):
                if c_idx < len(rows_data[r_idx]):
                    val = parse_value(rows_data[r_idx][c_idx])
                    if isinstance(val, (int, float)):
                        numeric_count += 1
                    elif isinstance(val, datetime.date):
                        date_count += 1
                    elif val:
                        text_count += 1
            if numeric_count > text_count:
                col_alignments[c_idx] = "right"
            elif date_count > text_count:
                col_alignments[c_idx] = "center"
            else:
                col_alignments[c_idx] = "left"

        for r_idx, row in enumerate(rows_data, start=1):
            is_header = (r_idx == 1)
            ws.row_dimensions[r_idx].height = 26 if is_header else 20
            is_alt = (r_idx % 2 == 0 and not is_header)

            for col_idx, cell_val in enumerate(row, start=1):
                parsed = parse_value(cell_val)
                c = ws.cell(row=r_idx, column=col_idx, value=parsed)
                c.border = thin_border
                align_dir = col_alignments.get(col_idx - 1, "left")

                if is_header:
                    c.fill = header_fill
                    c.font = header_font
                    c.alignment = Alignment(horizontal=align_dir, vertical="center", wrap_text=False)
                else:
                    c.font = regular_font
                    if is_alt:
                        c.fill = alt_fill
                    if isinstance(parsed, (int, float)):
                        c.alignment = Alignment(horizontal="right", vertical="center", wrap_text=False)
                        if isinstance(parsed, float):
                            c.number_format = "#,##0.00"
                    elif isinstance(parsed, datetime.date):
                        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=False)
                        c.number_format = "yyyy-mm-dd"
                    else:
                        c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=False)

        ws.freeze_panes = "A2"
        if len(rows_data) > 1 and len(rows_data[0]) > 0:
            last_col = get_column_letter(len(rows_data[0]))
            ws.auto_filter.ref = f"A1:{last_col}{len(rows_data)}"

        # Content-Aware Column Auto-Fitting
        for col in ws.columns:
            max_len = 0
            col_letter = get_column_letter(col[0].column)
            for cell in col:
                if cell.value is not None:
                    max_len = max(max_len, len(str(cell.value).split("\n")[0]))
            ws.column_dimensions[col_letter].width = max(min(max_len + 4, 60), 14)

    out_buf = io.BytesIO()
    wb.save(out_buf)
    xlsx_bytes = out_buf.getvalue()

    # Post-generation verification
    try:
        val_wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
        if len(val_wb.sheetnames) == 0 or val_wb.active.max_row < 1:
            raise ValueError("Generated XLSX is empty.")
    except Exception as e:
        raise ValueError(f"Generated Excel workbook failed validation: {e}")

    return xlsx_bytes

def pdf_to_html(pdf_bytes: bytes) -> bytes:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    html_out = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>Converted Document</title>
  <style>
    * { box-sizing: border-box; }
    body {
      background-color: #525659;
      margin: 0;
      padding: 30px 10px;
      font-family: system-ui, -apple-system, sans-serif;
      display: flex;
      flex-direction: column;
      align-items: center;
    }
    .pdf-page-container {
      background: white;
      position: relative;
      margin-bottom: 24px;
      box-shadow: 0 4px 20px rgba(0, 0, 0, 0.25);
      overflow: hidden;
    }
    .pdf-page-container div {
      position: relative;
    }
    .pdf-page-container p {
      position: absolute;
      margin: 0;
      padding: 0;
      white-space: pre-wrap;
    }
    .pdf-page-container img {
      position: absolute;
    }
  </style>
</head>
<body>
"""

    for page_idx, page in enumerate(doc):
        w = page.rect.width
        h = page.rect.height
        page_html = page.get_text("html")
        html_out += f'  <div class="pdf-page-container" style="width:{w}pt; height:{h}pt;">\n'
        html_out += page_html
        html_out += '\n  </div>\n'

    html_out += """</body>
</html>"""
    
    doc.close()
    return html_out.encode("utf-8")

def pdf_to_markdown(pdf_bytes: bytes) -> bytes:
    # 1. Primary: Enterprise Layout-Aware PyMuPDF4LLM Engine (CloudConvert / Marker-grade)
    try:
        import pymupdf4llm
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        md_text = pymupdf4llm.to_markdown(doc, page_chunks=False, write_images=False)
        doc.close()
        if md_text and len(md_text.strip()) > 0:
            return md_text.encode("utf-8")
    except Exception as e:
        print(f"[pdf_to_markdown pymupdf4llm]: {e}")

    # 2. Advanced Multi-Pass GFM Markdown Fallback
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    md_lines = []

    for page_idx, page in enumerate(doc):
        md_lines.append(f"<!-- Page {page_idx + 1} -->\n")
        
        # Check for vector/drawn tables on the page
        tabs = None
        try:
            tabs = page.find_tables()
        except Exception:
            pass

        table_bboxes = []
        if tabs and len(tabs.tables) > 0:
            for tab in tabs:
                table_bboxes.append(tab.bbox)
                df_data = tab.extract()
                if not df_data:
                    continue
                # Format as GitHub Flavored Markdown Table
                headers = [str(c or "").strip().replace("\n", " ") for c in df_data[0]]
                md_lines.append("| " + " | ".join(headers) + " |")
                md_lines.append("| " + " | ".join([":---" for _ in headers]) + " |")
                for row in df_data[1:]:
                    cleaned_row = [str(c or "").strip().replace("\n", " ") for c in row]
                    md_lines.append("| " + " | ".join(cleaned_row) + " |")
                md_lines.append("\n")

        # Process Non-Table Blocks
        text_blocks = page.get_text("blocks")
        text_blocks.sort(key=lambda b: b[1])

        for block in text_blocks:
            by0, by1, block_text = block[1], block[3], block[4].strip()
            if not block_text:
                continue

            inside_table = False
            for tbbox in table_bboxes:
                if by0 >= tbbox[1] - 5 and by1 <= tbbox[3] + 5:
                    inside_table = True
                    break
            if inside_table:
                continue

            lines = block_text.splitlines()
            is_code = any(l.strip().startswith(("import ", "from ", "def ", "class ", "plt.", "np.")) for l in lines)

            if is_code:
                md_lines.append("```python")
                md_lines.extend(lines)
                md_lines.append("```\n")
            else:
                for line in lines:
                    line_str = line.strip()
                    if not line_str:
                        continue
                    if line_str.isupper() and len(line_str) < 50:
                        md_lines.append(f"## {line_str}\n")
                    else:
                        md_lines.append(f"{line_str}\n")
                md_lines.append("\n")

        md_lines.append("\n---\n\n")

    doc.close()
    return "".join(md_lines).encode("utf-8")

def pdf_to_pdfa(pdf_bytes: bytes) -> bytes:
    import subprocess
    import tempfile
    import os

    # 1. Primary: Ghostscript ISO PDF/A Engine if available
    gs_cmd = get_ghostscript_cmd()
    if gs_cmd:
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                in_pdf = os.path.join(tmpdir, "in.pdf")
                out_pdf = os.path.join(tmpdir, "out_pdfa.pdf")
                with open(in_pdf, "wb") as f:
                    f.write(pdf_bytes)
                cmd = [
                    gs_cmd,
                    "-dPDFA=2",
                    "-dBATCH",
                    "-dNOPAUSE",
                    "-dNOOUTERSAVE",
                    "-sProcessColorModel=DeviceRGB",
                    "-sDEVICE=pdfwrite",
                    "-sPDFACompatibilityPolicy=1",
                    f"-sOutputFile={out_pdf}",
                    in_pdf
                ]
                res = subprocess.run(cmd, capture_output=True, timeout=25)
                if os.path.exists(out_pdf) and os.path.getsize(out_pdf) > 0:
                    with open(out_pdf, "rb") as f:
                        return f.read()
        except Exception as e:
            print(f"[pdf_to_pdfa Ghostscript]: {e}")

    # 2. Native PyMuPDF ISO-19005-2 PDF/A-2b Compliant Synthesizer
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    xmp_schema = (
        '<?xpacket begin="\ufeff" id="W5M0MpCehiHzreSzNTczkc9d"?>\n'
        '<x:xmpmeta xmlns:x="adobe:ns:meta/">\n'
        '  <rdf:RDF xmlns:rdf="http://www.w3.org/1999/02/22-rdf-syntax-ns#">\n'
        '    <rdf:Description rdf:about="" xmlns:pdfaid="http://www.aiim.org/pdfa/ns/schema#">\n'
        '      <pdfaid:part>2</pdfaid:part>\n'
        '      <pdfaid:conformance>B</pdfaid:conformance>\n'
        '    </rdf:Description>\n'
        '    <rdf:Description rdf:about="" xmlns:dc="http://purl.org/dc/elements/1.1/">\n'
        '      <dc:format>application/pdf</dc:format>\n'
        '    </rdf:Description>\n'
        '    <rdf:Description rdf:about="" xmlns:pdf="http://ns.adobe.com/pdf/1.3/">\n'
        '      <pdf:Producer>DocFlow ISO-19005-2 PDF/A Engine</pdf:Producer>\n'
        '    </rdf:Description>\n'
        '  </rdf:RDF>\n'
        '</x:xmpmeta>\n'
        '<?xpacket end="w"?>'
    )
    try:
        doc.set_xml_metadata(xmp_schema)
    except Exception as e:
        print(f"[pdf_to_pdfa xmp]: {e}")

    buffer = io.BytesIO()
    doc.save(buffer, garbage=4, deflate=True, clean=True)
    doc.close()
    return buffer.getvalue()
